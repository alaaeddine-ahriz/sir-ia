"""
Gestion des cours avec RAG - Système de génération de questions
Ce script permet de gérer des documents de cours par matière et de générer des questions de réflexion.
"""

# -----------------------------------------
# Imports et configuration d'environnement
# -----------------------------------------
import os
import time
import uuid
import warnings
import glob
import hashlib
import json
import tempfile
import io
from datetime import datetime
from pathlib import Path
from dotenv import load_dotenv

# Chargement des variables d'environnement
load_dotenv()

# Essai d'import de PyPDF2 au démarrage
PYPDF2_AVAILABLE = False
try:
    import PyPDF2
    PYPDF2_AVAILABLE = True
    print(f"PyPDF2 importé avec succès! Version: {PyPDF2.__version__}")
except ImportError as e:
    print(f"AVERTISSEMENT: PyPDF2 n'est pas disponible: {e}")
    print("L'extraction de texte PDF sera limitée.")

# Désactivation des avertissements LangSmith
os.environ["LANGCHAIN_TRACING_V2"] = "false"
os.environ["LANGCHAIN_TRACING"] = "false"
os.environ["LANGCHAIN_API_KEY"] = "not-needed"
os.environ["LANGCHAIN_PROJECT"] = "not-needed"
warnings.filterwarnings("ignore", category=Warning, module="langsmith")

# Imports pour les différents formats de documents
try:
    # PDF
    import pdfplumber
    # Word (docx)
    import docx
    # PowerPoint (pptx)
    from pptx import Presentation
    # OpenDocument (odt, odp)
    import odf.opendocument
    from odf.text import P
    from odf.teletype import extractText
except ImportError:
    print("Avertissement: Certaines bibliothèques de traitement de documents ne sont pas installées.")
    print("Exécutez 'pip install pdfplumber python-docx python-pptx odfpy PyPDF2' pour une prise en charge complète.")

# Imports pour le découpage de texte
from langchain_text_splitters import MarkdownHeaderTextSplitter, RecursiveCharacterTextSplitter

# Imports pour la base de données vectorielle
from pinecone import Pinecone, ServerlessSpec
from langchain_pinecone import PineconeEmbeddings, PineconeVectorStore

# Imports pour le système RAG
from langchain_openai import ChatOpenAI
from langchain.chains import create_retrieval_chain
from langchain.chains.combine_documents import create_stuff_documents_chain
from langchain.prompts import ChatPromptTemplate
from langchain import hub
from langchain.schema import Document

# -----------------------------------------
# Configuration du système
# -----------------------------------------

# Chemin vers le dossier des cours
COURS_DIR = os.path.join(os.getcwd(), "cours")
# Nom de l'index Pinecone
INDEX_NAME = "rag-sir"
# Fichier de suivi des mises à jour
METADATA_FILE = os.path.join(os.getcwd(), "metadata_cours.json")

# Dossier pour les sorties JSON
OUTPUTS_DIR = os.path.join(os.getcwd(), "outputs")

# -----------------------------------------
# Fonctions de gestion des fichiers
# -----------------------------------------
def initialiser_structure_dossiers():
    """
    Initialise la structure des dossiers pour les cours si elle n'existe pas.
    Crée un dossier par matière avec un fichier README explicatif.
    """
    # Créer le dossier principal des cours s'il n'existe pas
    if not os.path.exists(COURS_DIR):
        os.makedirs(COURS_DIR)
        print(f"Dossier principal des cours créé: {COURS_DIR}")
    
    # Liste des matières (à adapter selon vos besoins)
    matieres = ["SYD","TCP"]
    
    for matiere in matieres:
        matiere_dir = os.path.join(COURS_DIR, matiere)
        if not os.path.exists(matiere_dir):
            os.makedirs(matiere_dir)
            
            # Créer un fichier README explicatif
            readme_path = os.path.join(matiere_dir, "README.md")
            with open(readme_path, "w") as f:
                f.write(f"# Documents de cours pour {matiere}\n\n")
                f.write("Placez ici vos documents de cours pour cette matière.\n")
                f.write("Formats supportés: .md, .txt\n\n")
                f.write("Structure recommandée pour les fichiers markdown:\n")
                f.write("- Utilisez des titres ## pour les sections principales\n")
                f.write("- Chaque fichier doit traiter d'un concept ou d'une notion\n")
            
            print(f"Dossier pour la matière {matiere} créé avec un README explicatif")

def lire_fichiers_matiere(matiere):
    """
    Lit tous les fichiers de cours pour une matière donnée,
    avec prise en charge de formats variés (md, txt, pdf, docx, pptx, etc.).
    
    Args:
        matiere (str): Identifiant de la matière (ex: "SYD", "BD")
        
    Returns:
        list: Liste des documents avec leur contenu et métadonnées
    """
    matiere_dir = os.path.join(COURS_DIR, matiere)
    
    # Vérifier si le dossier existe
    if not os.path.exists(matiere_dir):
        print(f"Erreur: Le dossier de la matière {matiere} n'existe pas.")
        return []
    
    documents = []
    
    # Extensions supportées
    extensions = ["*.md", "*.txt", "*.pdf", "*.docx", "*.pptx", "*.doc", "*.odt", "*.odp"]
    
    # Vérifier si un dossier d'examens existe pour cette matière
    examens_dir = os.path.join(matiere_dir, "examens")
    exam_documents = []
    has_exam_folder = os.path.exists(examens_dir) and os.path.isdir(examens_dir)
    
    # Parcourir tous les fichiers avec les extensions supportées
    for ext in extensions:
        for file_path in glob.glob(os.path.join(matiere_dir, "**", ext), recursive=True):
            try:
                # Éviter de traiter les fichiers README
                if os.path.basename(file_path).lower() == "readme.md":
                    continue
                
                # Calculer le hash du fichier pour suivre les modifications
                file_hash = calculer_hash_fichier(file_path)
                
                # Extraire le contenu selon le type de fichier
                file_extension = os.path.splitext(file_path)[1].lower()
                content = extraire_contenu_fichier(file_path, file_extension)
                
                if not content or content.strip() == "":
                    print(f"Avertissement: Le fichier {file_path} semble vide après extraction.")
                    continue
                
                # Métadonnées du document
                relative_path = os.path.relpath(file_path, COURS_DIR)
                metadata = {
                    "source": relative_path,
                    "matiere": matiere,
                    "filename": os.path.basename(file_path),
                    "filetype": file_extension,
                    "file_hash": file_hash,
                    "updated_at": datetime.now().isoformat()
                }
                
                # Vérifier si le document est dans le dossier d'examens
                is_exam = "examens" in relative_path
                if is_exam:
                    metadata["is_exam"] = True
                    metadata["document_type"] = "exam"
                    exam_documents.append({"content": content, "metadata": metadata})
                else:
                    documents.append({"content": content, "metadata": metadata})
                
                print(f"Fichier lu: {relative_path}" + (" (examen)" if is_exam else ""))
                
            except Exception as e:
                print(f"Erreur lors de la lecture du fichier {file_path}: {e}")
    
    # Combiner les documents, plaçant les examens en premier pour leur donner plus de poids
    return exam_documents + documents

def calculer_hash_fichier(file_path):
    """
    Calcule un hash MD5 du contenu d'un fichier pour détecter les modifications.
    
    Args:
        file_path (str): Chemin vers le fichier
        
    Returns:
        str: Hash MD5 du fichier
    """
    hash_md5 = hashlib.md5()
    
    with open(file_path, "rb") as f:
        for chunk in iter(lambda: f.read(4096), b""):
            hash_md5.update(chunk)
            
    return hash_md5.hexdigest()

def extraire_contenu_fichier(file_path, file_extension):
    """
    Extrait le contenu textuel d'un fichier selon son format.
    
    Args:
        file_path (str): Chemin vers le fichier
        file_extension (str): Extension du fichier (avec le point)
        
    Returns:
        str: Contenu textuel du fichier
    """
    try:
        # Fichiers texte et markdown
        if file_extension in ['.txt', '.md']:
            with open(file_path, 'r', encoding='utf-8') as f:
                return f.read()
        
        # Fichiers PDF
        elif file_extension == '.pdf':
            text = ""
            try:
                with pdfplumber.open(file_path) as pdf:
                    for page in pdf.pages:
                        page_text = page.extract_text() or ""
                        text += page_text + "\n\n"
                return text
            except Exception as e:
                print(f"Erreur avec pdfplumber: {e}. Tentative alternative avec PyPDF2...")
                
                # Vérifier explicitement si PyPDF2 est disponible
                global PYPDF2_AVAILABLE
                if not PYPDF2_AVAILABLE:
                    try:
                        # Tenter d'importer à nouveau au cas où
                        import PyPDF2
                        PYPDF2_AVAILABLE = True
                        print("PyPDF2 importé avec succès pendant l'extraction!")
                    except ImportError as imp_err:
                        print(f"Échec de l'import PyPDF2: {imp_err}")
                        return f"[Erreur d'extraction PDF: Module PyPDF2 manquant, veuillez installer 'pip install PyPDF2==3.0.1']"
                
                # Maintenant PyPDF2 devrait être disponible
                try:
                    with open(file_path, 'rb') as file:
                        reader = PyPDF2.PdfReader(file)
                        text = ""
                        for page in reader.pages:
                            text += page.extract_text() + "\n\n"
                    return text
                except Exception as pdf_err:
                    print(f"Erreur avec PyPDF2: {pdf_err}")
                    return f"[Erreur d'extraction PDF avec PyPDF2: {str(pdf_err)}]"
        
        # Fichiers Word (DOCX)
        elif file_extension == '.docx':
            text = ""
            doc = docx.Document(file_path)
            for para in doc.paragraphs:
                text += para.text + "\n"
                # Tenter d'identifier les titres potentiels basés sur le style
                if para.style.name.startswith('Heading'):
                    heading_level = int(para.style.name[-1]) if para.style.name[-1].isdigit() else 2
                    prefix = '#' * heading_level
                    text = text[:-1] + f"\n{prefix} {para.text}\n"
            return text
        
        # Fichiers PowerPoint (PPTX)
        elif file_extension == '.pptx':
            text = ""
            pres = Presentation(file_path)
            for i, slide in enumerate(pres.slides):
                text += f"## Slide {i+1}\n\n"
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        text += shape.text + "\n"
                text += "\n"
            return text
        
        # Fichiers OpenDocument Text (ODT)
        elif file_extension == '.odt':
            textdoc = odf.opendocument.load(file_path)
            allparas = textdoc.getElementsByType(P)
            text = ""
            for para in allparas:
                text += extractText(para) + "\n"
            return text
        
        # Fichiers OpenDocument Presentation (ODP)
        elif file_extension == '.odp':
            doc = odf.opendocument.load(file_path)
            text = ""
            slide_num = 1
            
            # Récupérer tous les éléments de texte
            for para in doc.getElementsByType(P):
                content = extractText(para)
                if content.strip():
                    if "Slide" not in text[-20:] and len(text) > 0:
                        text += f"\n## Slide {slide_num}\n\n"
                        slide_num += 1
                    text += content + "\n"
            
            return text
        
        # Fichiers DOC (ancien format Word)
        elif file_extension == '.doc':
            # Essayer d'utiliser textract s'il est installé
            try:
                import textract
                text = textract.process(file_path).decode('utf-8')
                return text
            except ImportError:
                print(f"Erreur: textract non installé pour lire les fichiers .doc")
                print(f"Installez-le avec 'pip install textract'")
                return f"[Contenu du fichier .doc non extrait: {os.path.basename(file_path)}]"
        
        # Format non pris en charge
        else:
            print(f"Format de fichier non pris en charge: {file_extension}")
            return f"[Format non pris en charge: {file_extension}]"
    
    except Exception as e:
        print(f"Erreur lors de l'extraction du contenu de {file_path}: {e}")
        return f"[Erreur d'extraction: {str(e)}]"

def charger_metadata():
    """
    Charge les métadonnées des précédentes mises à jour depuis un fichier JSON.
    
    Returns:
        dict: Métadonnées des documents indexés
    """
    if os.path.exists(METADATA_FILE):
        try:
            with open(METADATA_FILE, 'r', encoding='utf-8') as f:
                return json.load(f)
        except Exception as e:
            print(f"Erreur lors du chargement des métadonnées: {e}")
            return {"matieres": {}}
    else:
        return {"matieres": {}}

def sauvegarder_metadata(metadata):
    """
    Sauvegarde les métadonnées des documents indexés dans un fichier JSON.
    
    Args:
        metadata (dict): Métadonnées à sauvegarder
    """
    try:
        with open(METADATA_FILE, 'w', encoding='utf-8') as f:
            json.dump(metadata, f, indent=2)
    except Exception as e:
        print(f"Erreur lors de la sauvegarde des métadonnées: {e}")

# -----------------------------------------
# Fonctions de gestion des fichiers de sortie
# -----------------------------------------
def creer_dossier_sortie():
    """
    Crée un dossier horodaté pour les sorties JSON.
    
    Returns:
        str: Chemin vers le dossier créé
    """
    # Créer le dossier principal de sorties s'il n'existe pas
    if not os.path.exists(OUTPUTS_DIR):
        os.makedirs(OUTPUTS_DIR)
    
    # Créer un sous-dossier horodaté
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    output_dir = os.path.join(OUTPUTS_DIR, f"outputs_{timestamp}")
    
    if not os.path.exists(output_dir):
        os.makedirs(output_dir)
        print(f"Dossier de sortie créé: {output_dir}")
    
    return output_dir

def sauvegarder_json(json_data, prefix, matiere, output_dir=None):
    """
    Sauvegarde un objet JSON dans un fichier horodaté.
    
    Args:
        json_data (str ou dict): Données JSON à sauvegarder
        prefix (str): Préfixe pour le nom du fichier
        matiere (str): Identifiant de la matière
        output_dir (str, optional): Dossier de sortie. Si None, utilise le dossier par défaut.
        
    Returns:
        str: Chemin vers le fichier créé
    """
    # Créer le dossier de sortie si nécessaire
    if output_dir is None:
        output_dir = creer_dossier_sortie()
    
    # Générer un nom de fichier unique
    timestamp = datetime.now().strftime("%H%M%S")
    filename = f"{prefix}_{matiere}_{timestamp}.json"
    filepath = os.path.join(output_dir, filename)
    
    # Convertir en chaîne JSON si nécessaire
    if isinstance(json_data, dict):
        json_str = json.dumps(json_data, ensure_ascii=False, indent=2)
    else:
        json_str = json_data
    
    # Écrire le fichier
    with open(filepath, 'w', encoding='utf-8') as f:
        f.write(json_str)
    
    print(f"Résultat JSON sauvegardé: {filepath}")
    return filepath

# -----------------------------------------
# Fonctions de traitement du texte
# -----------------------------------------
def split_document(document):
    """
    Divise un document en sections.
    Utilise des stratégies différentes selon le type de fichier et la présence d'en-têtes.
    
    Args:
        document (dict): Document avec contenu et métadonnées
        
    Returns:
        list: Liste des sections avec leurs métadonnées
    """
    content = document["content"]
    metadata = document["metadata"]
    
    # Vérifier la présence d'en-têtes markdown (##, ###) dans le contenu
    has_markdown_headers = '##' in content
    
    # Si le document est un markdown ou contient des en-têtes, utiliser la méthode par en-têtes
    if metadata["filetype"] == ".md" or has_markdown_headers:
        headers_to_split_on = [
            ("##", "Header 2"),
            ("###", "Header 3")
        ]
        
        try:
            markdown_splitter = MarkdownHeaderTextSplitter(
                headers_to_split_on=headers_to_split_on, strip_headers=False
            )
            
            splits = markdown_splitter.split_text(content)
            
            # Si aucun en-tête n'a été trouvé, utiliser la méthode par caractères
            if not splits:
                return split_by_characters(content, metadata)
                
            # Ajouter les métadonnées du document à chaque split
            for split in splits:
                split.metadata.update(metadata)
                
            return splits
            
        except Exception as e:
            print(f"Erreur lors du découpage markdown: {e}")
            # Fallback sur le découpage par caractères
            return split_by_characters(content, metadata)
    
    # Pour les autres types de fichiers ou documents sans en-têtes: découpage par caractères et paragraphes
    else:
        return split_by_paragraphs(content, metadata)

def split_by_characters(content, metadata):
    """
    Divise un document en chunks par caractères (méthode de secours).
    
    Args:
        content (str): Contenu du document
        metadata (dict): Métadonnées du document
        
    Returns:
        list: Liste des sections avec leurs métadonnées
    """
    text_splitter = RecursiveCharacterTextSplitter(
        chunk_size=1000,
        chunk_overlap=200,
        length_function=len,
    )
    
    docs = text_splitter.create_documents([content], [metadata])
    return docs

def split_by_paragraphs(content, metadata):
    """
    Divise un document en chunks par paragraphes,
    plus adapté pour les documents non structurés.
    
    Args:
        content (str): Contenu du document
        metadata (dict): Métadonnées du document
        
    Returns:
        list: Liste des sections avec leurs métadonnées
    """
    # Découpage par paragraphes avec chevauchement
    text_splitter = RecursiveCharacterTextSplitter(
        chunk_size=1000,
        chunk_overlap=250,
        separators=["\n\n", "\n", ". ", " ", ""],
        length_function=len,
    )
    
    docs = text_splitter.create_documents([content], [metadata])
    return docs

# -----------------------------------------
# Configuration de la base de données vectorielle
# -----------------------------------------
def initialize_pinecone():
    """
    Initialise la connexion à Pinecone et crée un index si nécessaire.
    
    Returns:
        tuple: (client Pinecone, nom de l'index, spécifications)
    """
    pinecone_api_key = os.environ.get("PINECONE_API_KEY")
    cloud = os.environ.get('PINECONE_CLOUD') or 'aws'
    region = os.environ.get('PINECONE_REGION') or 'us-east-1'
    
    pc = Pinecone(api_key=pinecone_api_key)
    spec = ServerlessSpec(cloud=cloud, region=region)
    
    return pc, INDEX_NAME, spec

def setup_embeddings():
    """
    Initialise le modèle d'embedding pour la vectorisation du texte.
    
    Returns:
        PineconeEmbeddings: Instance configurée du modèle d'embedding
    """
    model_name = 'multilingual-e5-large'
    
    embeddings = PineconeEmbeddings(
        model=model_name,
        pinecone_api_key=os.environ.get('PINECONE_API_KEY')
    )
    
    return embeddings

def create_or_get_index(pc, index_name, embeddings, spec):
    """
    Crée un nouvel index si nécessaire ou récupère un index existant.
    
    Args:
        pc: Client Pinecone
        index_name (str): Nom de l'index
        embeddings: Modèle d'embedding
        spec: Spécifications de l'index
        
    Returns:
        object: Index Pinecone
    """
    if index_name not in pc.list_indexes().names():
        print(f"Création d'un nouvel index: {index_name}")
        pc.create_index(
            name=index_name,
            dimension=embeddings.dimension,
            metric="cosine",
            spec=spec
        )
        # Attendre que l'index s'initialise
        print("Attente de l'initialisation de l'index...")
        time.sleep(10)
    
    # Afficher l'état de l'index
    print("État actuel de l'index:")
    index = pc.Index(index_name)
    index_stats = index.describe_index_stats()
    
    # Utiliser un dictionnaire au lieu d'un objet spécifique de Pinecone
    stats_dict = {
        "dimension": index_stats.dimension,
        "index_fullness": index_stats.index_fullness,
        "namespaces": {},
        "total_vector_count": index_stats.total_vector_count
    }
    
    # Ajouter les informations sur les namespaces
    if hasattr(index_stats, 'namespaces') and index_stats.namespaces:
        for ns_name, ns_data in index_stats.namespaces.items():
            stats_dict["namespaces"][ns_name] = {"vector_count": ns_data.vector_count}
    
    print(json.dumps(stats_dict, indent=2))
    print("\n")
    
    return index

def get_matiere_namespace(matiere):
    """
    Génère un namespace standardisé pour une matière.
    
    Args:
        matiere (str): Identifiant de la matière
        
    Returns:
        str: Namespace standardisé
    """
    return f"matiere-{matiere.lower()}"

def upsert_documents(pc, index_name, embeddings, matiere, docs):
    """
    Met à jour l'index vectoriel avec de nouveaux documents ou des documents modifiés.
    
    Args:
        pc: Client Pinecone
        index_name (str): Nom de l'index
        embeddings: Modèle d'embedding
        matiere (str): Identifiant de la matière
        docs (list): Liste des sections de documents
        
    Returns:
        tuple: (store vectoriel, namespace utilisé)
    """
    if not docs:
        print(f"Aucun document à insérer pour la matière {matiere}")
        return None, None
    
    namespace = get_matiere_namespace(matiere)
    
    try:
        print(f"Insertion de {len(docs)} sections de documents dans l'espace '{namespace}'")
        vector_store = PineconeVectorStore.from_documents(
            documents=docs,
            index_name=index_name,
            embedding=embeddings,
            namespace=namespace
        )
        print("Insertion réussie")
        return vector_store, namespace
    except Exception as e:
        print(f"Erreur lors de l'insertion: {e}")
        return None, None

def delete_documents(pc, index_name, matiere, file_paths):
    """
    Supprime les documents d'un index vectoriel pour les fichiers qui ont été supprimés,
    en utilisant une approche en deux étapes: d'abord identifier les vecteurs par requête,
    puis les supprimer par ID.
    
    Args:
        pc: Client Pinecone
        index_name (str): Nom de l'index
        matiere (str): Identifiant de la matière
        file_paths (list): Liste des chemins de fichiers supprimés
        
    Returns:
        bool: True si des suppressions ont été effectuées, False sinon
    """
    if not file_paths:
        return False
    
    namespace = get_matiere_namespace(matiere)
    
    try:
        index = pc.Index(index_name)
        deleted_count = 0
        
        # Afficher les informations de l'index pour le débogage
        try:
            index_stats = index.describe_index_stats()
            print(f"Statistiques de l'index '{index_name}':")
            print(f"- Dimension: {index_stats.dimension}")
            print(f"- Total vecteurs: {index_stats.total_vector_count}")
            
            # Afficher les informations du namespace si disponible
            if hasattr(index_stats, 'namespaces') and namespace in index_stats.namespaces:
                ns_stats = index_stats.namespaces[namespace]
                print(f"- Namespace '{namespace}': {ns_stats.vector_count} vecteurs")
                
            # Obtenir la dimension du vecteur pour la requête
            dimension = index_stats.dimension
        except Exception as stats_error:
            print(f"Impossible d'obtenir les statistiques de l'index: {stats_error}")
            # Valeur par défaut pour la dimension
            dimension = 1024
        
        # Traiter chaque fichier à supprimer
        for file_path in file_paths:
            print(f"\nSuppression des vecteurs pour: {file_path}")
            
            # Étape 1: Identifier les vecteurs par requête
            try:
                # Créer un vecteur zéro pour la requête
                zero_vector = [0.0] * dimension
                
                # Exécuter une requête avec un grand nombre de résultats
                print("1. Recherche des vecteurs à supprimer...")
                query_results = index.query(
                    namespace=namespace,
                    vector=zero_vector,
                    top_k=1000,  # Augmenter si nécessaire
                    include_metadata=True
                )
                
                # Préparer la liste des IDs à supprimer
                ids_to_delete = []
                
                # Vérifier les résultats de la requête
                if hasattr(query_results, 'matches') and query_results.matches:
                    print(f"   Vecteurs trouvés: {len(query_results.matches)}")
                    
                    # Identifier les vecteurs correspondant au fichier supprimé
                    for match in query_results.matches:
                        # Vérifier si le vecteur a des métadonnées
                        if hasattr(match, 'metadata') and match.metadata:
                            # Vérifier si le champ source correspond exactement au fichier
                            if 'source' in match.metadata and match.metadata['source'] == file_path:
                                ids_to_delete.append(match.id)
                                # Afficher un exemple pour confirmation
                                if len(ids_to_delete) == 1:
                                    print(f"   Exemple de métadonnées trouvées: {match.metadata}")
                
                # Étape 2: Supprimer les vecteurs identifiés
                if ids_to_delete:
                    print(f"2. Suppression de {len(ids_to_delete)} vecteurs...")
                    
                    # Supprimer par lots pour éviter les limitations d'API
                    batch_size = 100
                    for i in range(0, len(ids_to_delete), batch_size):
                        batch = ids_to_delete[i:i+batch_size]
                        delete_result = index.delete(
                            ids=batch,
                            namespace=namespace
                        )
                        
                        if hasattr(delete_result, 'deleted_count'):
                            print(f"   Lot {i//batch_size + 1}: {delete_result.deleted_count} vecteurs supprimés")
                        else:
                            print(f"   Lot {i//batch_size + 1}: suppression effectuée")
                    
                    deleted_count += len(ids_to_delete)
                    print(f"✅ {len(ids_to_delete)} vecteurs supprimés pour {file_path}")
                else:
                    print(f"❌ Aucun vecteur trouvé avec source={file_path}")
                    
                    # Afficher un échantillon de métadonnées pour le débogage
                    if hasattr(query_results, 'matches') and query_results.matches:
                        print("Exemples de métadonnées dans l'index (pour déboguer):")
                        samples_shown = 0
                        for match in query_results.matches:
                            if hasattr(match, 'metadata') and match.metadata and samples_shown < 3:
                                print(f"- ID: {match.id}")
                                for key, value in match.metadata.items():
                                    print(f"  {key}: {value}")
                                samples_shown += 1
                                print()
                
            except Exception as e:
                print(f"Erreur lors de la recherche/suppression pour {file_path}: {e}")
        
        # Résumer les résultats
        if deleted_count > 0:
            print(f"\n✅ Total: {deleted_count} vecteurs supprimés pour {len(file_paths)} fichiers")
            return True
        else:
            print(f"\n❌ Aucun vecteur n'a pu être supprimé pour les fichiers spécifiés")
            return False
            
    except Exception as e:
        print(f"Erreur générale lors de la suppression des documents: {e}")
        return False

# -----------------------------------------
# Gestion des mises à jour
# -----------------------------------------
def mettre_a_jour_matiere(pc, index_name, embeddings, matiere):
    """
    Met à jour l'index vectoriel pour une matière spécifique,
    en ne traitant que les fichiers nouveaux ou modifiés et en supprimant les fichiers disparus.
    
    Args:
        pc: Client Pinecone
        index_name (str): Nom de l'index
        embeddings: Modèle d'embedding
        matiere (str): Identifiant de la matière
        
    Returns:
        bool: True si des mises à jour ont été effectuées, False sinon
    """
    # Charger les métadonnées des précédentes indexations
    metadata = charger_metadata()
    
    # Initialiser la section pour cette matière si elle n'existe pas
    if matiere not in metadata["matieres"]:
        metadata["matieres"][matiere] = {
            "fichiers": {},
            "derniere_mise_a_jour": None
        }
    
    # Récupérer les fichiers précédemment indexés
    fichiers_matieres = metadata["matieres"][matiere]["fichiers"]
    
    # Lire tous les fichiers de cours actuels pour cette matière
    documents = lire_fichiers_matiere(matiere)
    sources_actuelles = set()
    if documents:
        sources_actuelles = {doc["metadata"]["source"] for doc in documents}
    
    # Déterminer quels fichiers ont été supprimés
    sources_indexees = set(fichiers_matieres.keys())
    fichiers_supprimes = sources_indexees - sources_actuelles
    
    # Afficher les fichiers supprimés
    if fichiers_supprimes:
        print(f"Fichiers supprimés détectés: {len(fichiers_supprimes)}")
        for source in fichiers_supprimes:
            print(f"  - {source}")
        
        # Supprimer les fichiers de l'index
        delete_documents(pc, index_name, matiere, list(fichiers_supprimes))
        
        # Mettre à jour les métadonnées pour refléter les suppressions
        for source in fichiers_supprimes:
            if source in fichiers_matieres:
                del fichiers_matieres[source]
    
    if not documents:
        print(f"Aucun document trouvé pour la matière {matiere}")
        # Sauvegarder les métadonnées si des fichiers ont été supprimés
        if fichiers_supprimes:
            metadata["matieres"][matiere]["derniere_mise_a_jour"] = datetime.now().isoformat()
            sauvegarder_metadata(metadata)
            return True
        return False
    
    # Déterminer quels fichiers ont été modifiés ou ajoutés
    docs_a_traiter = []
    
    for doc in documents:
        source = doc["metadata"]["source"]
        file_hash = doc["metadata"]["file_hash"]
        
        # Vérifier si le fichier est nouveau ou modifié
        if source not in fichiers_matieres or fichiers_matieres[source] != file_hash:
            print(f"Fichier nouveau ou modifié détecté: {source}")
            docs_a_traiter.append(doc)
            # Mettre à jour le hash dans les métadonnées
            fichiers_matieres[source] = file_hash
    
    # Si aucun document n'a été modifié ni supprimé, terminer
    if not docs_a_traiter and not fichiers_supprimes:
        print(f"Aucune mise à jour nécessaire pour la matière {matiere}")
        return False
    
    # Si des documents ont été modifiés ou ajoutés, les traiter
    if docs_a_traiter:
        # Diviser les documents en sections
        sections = []
        for doc in docs_a_traiter:
            doc_sections = split_document(doc)
            sections.extend(doc_sections)
        
        print(f"Nombre total de sections à indexer: {len(sections)}")
        
        # Mettre à jour l'index vectoriel
        vector_store, namespace = upsert_documents(pc, index_name, embeddings, matiere, sections)
        
        if not vector_store and not fichiers_supprimes:
            return False
    
    # Mettre à jour les métadonnées
    metadata["matieres"][matiere]["derniere_mise_a_jour"] = datetime.now().isoformat()
    sauvegarder_metadata(metadata)
    return True

# -----------------------------------------
# Configuration du système RAG
# -----------------------------------------
def setup_rag_system(index_name, embeddings, matiere, custom_prompt=None, output_format="text"):
    """
    Configure le système RAG pour une matière spécifique.
    
    Args:
        index_name (str): Nom de l'index Pinecone
        embeddings: Modèle d'embedding
        matiere (str): Identifiant de la matière
        custom_prompt: Prompt personnalisé facultatif
        output_format (str): Format de sortie ("text" ou "json")
        
    Returns:
        object: Chaîne de récupération configurée
    """
    namespace = get_matiere_namespace(matiere)
    
    # Créer un store vectoriel pour cette matière
    vector_store = PineconeVectorStore(
        index_name=index_name,
        embedding=embeddings,
        namespace=namespace
    )
    
    # Configurer le prompt de question-réponse
    if custom_prompt:
        retrieval_qa_chat_prompt = custom_prompt
    elif output_format == "json":
        retrieval_qa_chat_prompt = creer_prompt_json(matiere)
    else:
        retrieval_qa_chat_prompt = hub.pull("langchain-ai/retrieval-qa-chat")
    
    # Configurer le retriever
    retriever = vector_store.as_retriever()
    
    # Configurer le modèle de langage
    llm = ChatOpenAI(
        openai_api_key=os.environ.get('OPENAI_API_KEY'),
        model_name='gpt-4o-mini',
        temperature=0.0
    )
    
    # Créer la chaîne de traitement des documents
    combine_docs_chain = create_stuff_documents_chain(llm, retrieval_qa_chat_prompt)
    
    # Créer la chaîne de récupération
    retrieval_chain = create_retrieval_chain(retriever, combine_docs_chain)
    
    return retrieval_chain

def creer_prompt_json(matiere):
    """
    Crée un prompt personnalisé pour générer des réponses au format JSON sans les sources.
    Les sources seront ajoutées programmatiquement à partir des documents réels.
    
    Args:
        matiere (str): Identifiant de la matière
    
    Returns:
        ChatPromptTemplate: Template de prompt pour générer du JSON
    """
    TEMPLATE_JSON = """
    Vous êtes un assistant pédagogique IA spécialisé dans la matière {matiere}, disposant d'un accès direct aux documents de cours via un système de recherche sémantique (RAG).

    Sur la base des extraits de cours suivants:
    {context}
    
    Répondez à cette question: {input}
    
    Votre réponse DOIT être au format JSON suivant, SANS les sources (j'ajouterai les sources moi-même):
    
    ```json
    {{
        "réponse": "La réponse complète à la question",
        "niveau_confiance": 0.95
    }}
    ```
    
    CHAMPS OPTIONNELS:
    - Si pertinent, vous pouvez ajouter un champ "concepts_clés": ["concept1", "concept2", "concept3"]
    - Ce champ est OPTIONNEL et ne doit être inclus que si vous pouvez identifier clairement des concepts clés
    
    IMPORTANT:
    - Ne mentionnez pas et n'ajoutez pas de sources dans votre réponse
    - Concentrez-vous principalement sur la réponse
    - N'incluez que les champs spécifiés ci-dessus
    
    Ne répondez qu'avec ce format JSON, sans aucun texte avant ou après.
    """
    
    return ChatPromptTemplate.from_template(TEMPLATE_JSON).partial(matiere=matiere)

def creer_prompt_tuteur(matiere, output_format="text"):
    """
    Crée un prompt personnalisé pour le tuteur d'une matière spécifique.
    
    Args:
        matiere (str): Identifiant de la matière
        output_format (str): Format de sortie ("text" ou "json")
    
    Returns:
        ChatPromptTemplate: Template de prompt pour le tuteur
    """
    if output_format == "json":
        TEMPLATE_TUTEUR_JSON = """
        Vous êtes un tuteur IA spécialisé dans la matière {matiere}, disposant d'un accès direct aux documents de cours via un système de recherche sémantique (RAG).

        Votre tâche est de générer une question de réflexion sur le concept demandé, en vous basant strictement sur les extraits suivants:

        {context}

        IMPORTANT CONCERNANT LES EXAMENS:
        Certains des documents fournis peuvent être des fichiers d'examens (identifiés par "is_exam": true dans les métadonnées).
        Si ces documents sont présents parmi les sources, vous devez:
        1. Observer attentivement le style, le niveau et le type de questions posées par les professeurs dans ces examens
        2. Formuler votre question dans un style similaire, visant le même niveau de difficulté et de réflexion
        3. S'inspirer de la structure et de la formulation des questions d'examen tout en restant pertinent au concept demandé

        La question doit :
        1. Solliciter l'analyse critique d'un concept ou d'une relation entre plusieurs notions.
        2. Être formulée de manière claire, concise et précise, avec un vocabulaire académique adapté.
        3. Favoriser une réponse argumentée plutôt qu'une simple définition.
        4. Ressembler au style et au niveau des questions d'examen si des documents d'examen font partie des sources.

        Votre réponse DOIT être strictement au format JSON suivant, SANS les sources (j'ajouterai les sources moi-même):

        ```json
        {{
            "question": "La question de réflexion complète",
            "concepts_abordés": ["concept1", "concept2", "concept3"],
            "niveau_difficulté": "avancé", // "débutant", "intermédiaire" ou "avancé"
            "compétences_visées": ["analyse critique", "synthèse", "application pratique"],
            "éléments_réponse": [
                "Élément 1 attendu dans la réponse",
                "Élément 2 attendu dans la réponse",
                "Élément 3 attendu dans la réponse"
            ],
            "basé_sur_examen": true // ou false, selon si des documents d'examen ont influencé la formulation
        }}
        ```

        IMPORTANT:
        - Ne mentionnez pas et n'ajoutez pas de sources dans votre réponse
        - Concentrez-vous uniquement sur la question et les éléments associés
        - N'incluez aucun autre champ que ceux spécifiés ci-dessus
        - Les "éléments_réponse" doivent être des points clés qu'un étudiant devrait aborder dans sa réponse
        - Basez ces éléments uniquement sur le contenu des documents sources fournis

        Ne répondez qu'avec ce format JSON, sans aucun texte avant ou après.
        """
        
        return ChatPromptTemplate.from_template(TEMPLATE_TUTEUR_JSON).partial(matiere=matiere)
    else:
        TEMPLATE_TUTEUR = """
        Vous êtes un tuteur IA spécialisé dans la matière {matiere}, disposant d'un accès direct aux documents de cours via un système de recherche sémantique (RAG).

        Votre tâche est de générer une seule question de réflexion, en vous basant strictement sur les extraits suivants:

        {context}

        IMPORTANT CONCERNANT LES EXAMENS:
        Certains des documents fournis peuvent être des fichiers d'examens (identifiés par "is_exam": true dans les métadonnées).
        Si ces documents sont présents parmi les sources, vous devez:
        1. Observer attentivement le style, le niveau et le type de questions posées par les professeurs dans ces examens
        2. Formuler votre question dans un style similaire, visant le même niveau de difficulté et de réflexion
        3. S'inspirer de la structure et de la formulation des questions d'examen tout en restant pertinent au concept demandé

        La question doit :
        1. Solliciter l'analyse critique d'un concept ou d'une relation entre plusieurs notions.
        2. Être formulée de manière claire, concise et précise, avec un vocabulaire académique adapté.
        3. Favoriser une réponse argumentée plutôt qu'une simple définition.
        4. Ressembler au style et au niveau des questions d'examen si des documents d'examen font partie des sources.

        # Instructions

        - N'utilisez que les passages extraits des documents de cours ci-dessus.
        - Ne proposez qu'une question ouverte et directe en une ligne.
        - N'ajoutez ni explication ni sous-questions.
        - Les questions ne doivent pas demander d'analyse.
        - Référez-vous strictement aux extraits listés.

        Votre question: 
        """
        return ChatPromptTemplate.from_template(TEMPLATE_TUTEUR).partial(matiere=matiere)

def creer_prompt_evaluateur(matiere, output_format="json"):
    """
    Crée un prompt personnalisé pour l'évaluateur d'une réponse d'étudiant.
    
    Args:
        matiere (str): Identifiant de la matière
        output_format (str): Format de sortie (uniquement "json" supporté)
    
    Returns:
        ChatPromptTemplate: Template de prompt pour l'évaluateur
    """
    TEMPLATE_EVALUATEUR = """
    Vous êtes un examinateur académique automatisé spécialisé dans la matière {matiere}. 
    Votre rôle est d'évaluer la réponse d'un étudiant à une question de réflexion, en vous basant strictement sur le contenu du cours.

    Question posée: {question}
    
    Réponse de l'étudiant: {student_response}
    
    Contexte du cours (utilisez ces extraits comme référence pour évaluer la pertinence et l'exactitude des connaissances):
    {context}
    
    IMPORTANT CONCERNANT LES EXAMENS:
    Certains des documents fournis peuvent être des fichiers d'examens (identifiés par "is_exam": true dans les métadonnées).
    Si ces documents sont présents parmi les sources, vous devez:
    1. Observer attentivement le style et les critères d'évaluation utilisés par les professeurs dans ces examens
    2. Appliquer des standards d'évaluation similaires à ceux qu'un professeur utiliserait pour cette matière
    3. Tenir compte du niveau de difficulté et de précision attendu dans les examens officiels
    
    Procédez de façon rigoureuse, pédagogique et structurée selon les étapes suivantes:
    
    1. Évaluez la réponse en considérant:
       - Pertinence des idées: Les arguments répondent-ils à la question?
       - Qualité de l'argumentation: Les idées sont-elles développées et logiques?
       - Maîtrise des connaissances: L'étudiant utilise-t-il correctement les concepts du cours?
       - Originalité et pensée critique: La réponse montre-t-elle une réflexion personnelle?
       - Clarté et structure: L'expression est-elle compréhensible et organisée?
    
    2. Rédigez une réponse modèle concise mais complète
    
    3. Identifiez 3 points forts et 3 points à améliorer
    
    4. Attribuez une note sur 100 et justifiez-la
    
    5. Proposez un conseil personnalisé pour amélioration
    
    Votre évaluation DOIT être retournée strictement au format JSON suivant:
    
    ```json
    {
        "note": 85,
        "points_forts": [
            "Point fort 1",
            "Point fort 2",
            "Point fort 3"
        ],
        "points_ameliorer": [
            "Point à améliorer 1",
            "Point à améliorer 2",
            "Point à améliorer 3"
        ],
        "reponse_modele": "Une réponse modèle concise mais complète",
        "justification_note": "Explication détaillée de la note attribuée",
        "conseil_personnalise": "Un conseil spécifique pour aider l'étudiant à progresser",
        "basé_sur_examen": true
    }
    ```
    
    IMPORTANT:
    - Le champ "basé_sur_examen" doit être true si des documents d'examen ont influencé l'évaluation, false sinon
    - Ne mentionnez pas et n'ajoutez pas de sources dans votre réponse
    - N'incluez aucun autre champ que ceux spécifiés ci-dessus
    - Soyez rigoureux mais juste dans votre évaluation
    
    Ne répondez qu'avec ce format JSON, sans aucun texte avant ou après.
    """
    
    return ChatPromptTemplate.from_template(TEMPLATE_EVALUATEUR).partial(matiere=matiere)

# -----------------------------------------
# Fonctions d'interrogation
# -----------------------------------------
def interroger_matiere(index_name, embeddings, matiere, query, custom_prompt=None, output_format="text", save_output=True):
    """
    Interroge spécifiquement les documents d'une matière.
    
    Args:
        index_name (str): Nom de l'index Pinecone
        embeddings: Modèle d'embedding
        matiere (str): Identifiant de la matière
        query (str): Question de l'utilisateur
        custom_prompt: Prompt personnalisé facultatif
        output_format (str): Format de sortie ("text" ou "json")
        save_output (bool): Indique si la sortie JSON doit être sauvegardée
        
    Returns:
        dict: Réponse du système RAG
    """
    # Configurer le système RAG
    retrieval_chain = setup_rag_system(index_name, embeddings, matiere, custom_prompt, output_format)
    
    print(f"\n\n--- Requête pour {matiere}: '{query}' ---")
    
    # Exécuter la requête
    response = retrieval_chain.invoke({"input": query})
    
    # Dossier de sortie pour la session
    output_dir = None
    if save_output and output_format == "json":
        output_dir = creer_dossier_sortie()
    
    # Formater la réponse si le format de sortie est JSON
    if output_format == "json":
        try:
            import json
            import re
            
            # Extraire le JSON de la réponse
            json_answer = response['answer']
            
            # Rechercher le JSON entre les délimiteurs ```json et ```
            json_pattern = r"```(?:json)?(.*?)```"
            match = re.search(json_pattern, json_answer, re.DOTALL)
            
            if match:
                # Extraire le contenu JSON
                json_str = match.group(1).strip()
                response_json = json.loads(json_str)
            else:
                # Essayer de charger directement la réponse comme JSON
                json_str = json_answer.strip()
                response_json = json.loads(json_str)
            
            # Vérifier si des documents d'examen sont présents dans le contexte
            has_exam_docs = any(
                doc.metadata.get("is_exam", False) for doc in response.get("context", [])
            )
            
            # Ajouter l'information sur l'utilisation d'examens si nécessaire
            if "basé_sur_examen" not in response_json:
                response_json["basé_sur_examen"] = has_exam_docs
            
            # S'assurer que le champ "concepts_clés" existe (même vide) pour la cohérence des réponses
            if "concepts_clés" not in response_json:
                # Champ facultatif, nous ne l'ajoutons pas s'il n'est pas présent
                pass
            
            # Ajouter les sources à la réponse
            sources = []
            for i, doc in enumerate(response["context"]):
                source_entry = {
                    "document": i + 1,
                    "source": doc.metadata.get('source', 'Source inconnue'),
                    "is_exam": doc.metadata.get('is_exam', False)
                }
                
                # Ajouter la section si elle existe
                if "Header 2" in doc.metadata:
                    source_entry["section"] = doc.metadata["Header 2"]
                elif "Header 3" in doc.metadata:
                    source_entry["section"] = doc.metadata["Header 3"]
                
                # Limiter le contenu pour éviter des extraits trop longs
                max_content_length = 250  # Caractères
                content = doc.page_content
                if len(content) > max_content_length:
                    content = content[:max_content_length] + "..."
                
                source_entry["contenu"] = content
                sources.append(source_entry)
            
            # Ajouter les sources à la réponse JSON
            response_json["sources"] = sources
            
            # Ajouter la requête originale et des métadonnées
            response_json["requête_originale"] = query
            response_json["matière"] = matiere
            response_json["date_génération"] = datetime.now().isoformat()
            
            # Mettre à jour la réponse
            formatted_json = json.dumps(response_json, ensure_ascii=False, indent=2)
            
            # Stocker le JSON complet dans la réponse
            response['answer'] = formatted_json
            
            # Sauvegarder le résultat si demandé
            if save_output:
                prefix = "reponse"
                sauvegarder_json(formatted_json, prefix, matiere, output_dir)
            
            print("\nRéponse (format JSON):")
            print(formatted_json)
            
        except Exception as e:
            print(f"Erreur lors du traitement de la réponse JSON: {e}")
            print(f"Contenu de la réponse: {response['answer'][:100]}...")
    
    else:
        # Pour le format texte, afficher simplement la réponse
        print("\nRéponse:")
        print(response["answer"])
    
    return response

def generer_question_reflexion(index_name, embeddings, matiere, concept_cle=None, output_format="text", save_output=True):
    """
    Génère une question de réflexion sur un concept clé dans une matière spécifique.
    
    Args:
        index_name (str): Nom de l'index Pinecone
        embeddings: Modèle d'embedding
        matiere (str): Identifiant de la matière (ex: "SYD", "BD")
        concept_cle (str, optional): Concept sur lequel générer une question. Si None ou vide, une question générale sera générée.
        output_format (str): Format de sortie ("text" ou "json")
        save_output (bool): Indique si la sortie JSON doit être sauvegardée
        
    Returns:
        str: Question de réflexion générée
    """
    # Créer le prompt tuteur pour cette matière
    tuteur_prompt = creer_prompt_tuteur(matiere, output_format)
    
    # Construire la requête en fonction de si concept_cle est fourni ou non
    if concept_cle and concept_cle.strip():
        query = f"Générer une question de réflexion sur le concept: {concept_cle}"
    else:
        query = f"Générer une question de réflexion générale sur la matière"
    
    # Interroger la matière avec ce prompt
    result = interroger_matiere(
        index_name=index_name, 
        embeddings=embeddings,
        matiere=matiere, 
        query=query,
        custom_prompt=tuteur_prompt,
        output_format=output_format,
        save_output=save_output
    )
    
    # Si le format est JSON, vérifier si des examens ont été utilisés pour la génération
    if output_format == "json" and isinstance(result.get("context"), list):
        # Vérifier si des documents d'examen sont présents dans le contexte
        has_exam_docs = any(
            doc.metadata.get("is_exam", False) for doc in result.get("context", [])
        )
        
        # Ajouter cette information à la réponse JSON si elle n'est pas déjà présente
        try:
            import json
            import re
            
            # Extraire le JSON de la réponse
            json_answer = result['answer']
            
            # Rechercher le JSON entre les délimiteurs ```json et ```
            json_pattern = r"```(?:json)?(.*?)```"
            match = re.search(json_pattern, json_answer, re.DOTALL)
            
            if match:
                # Extraire le contenu JSON
                json_str = match.group(1).strip()
                response_json = json.loads(json_str)
            else:
                # Essayer de charger directement la réponse comme JSON
                json_str = json_answer.strip()
                response_json = json.loads(json_str)
            
            # Si le champ "basé_sur_examen" n'est pas déjà défini, l'ajouter
            if "basé_sur_examen" not in response_json:
                response_json["basé_sur_examen"] = has_exam_docs
                
                # Mettre à jour la réponse JSON
                formatted_json = json.dumps(response_json, ensure_ascii=False, indent=2)
                result['answer'] = formatted_json
                
        except Exception as e:
            print(f"Erreur lors du traitement du JSON pour les examens: {e}")
    
    return result["answer"]

def evaluer_reponse_etudiant(index_name, embeddings, matiere, question, student_response, output_format="json", save_output=True):
    """
    Évalue la réponse d'un étudiant à une question de réflexion,
    en se basant sur le contenu des cours pour la matière spécifique.
    
    Args:
        index_name (str): Nom de l'index Pinecone
        embeddings: Modèle d'embedding
        matiere (str): Identifiant de la matière (ex: "SYD", "BD")
        question (str): Question de réflexion posée
        student_response (str): Réponse de l'étudiant à évaluer
        output_format (str): Format de sortie (uniquement "json" supporté)
        save_output (bool): Indique si la sortie JSON doit être sauvegardée
        
    Returns:
        dict: Résultat de l'évaluation avec la note, les points forts/faibles, etc.
    """
    # Valider le format de sortie (seul json est supporté pour cette fonction)
    if output_format != "json":
        print("Avertissement: Seul le format 'json' est supporté pour l'évaluation. Utilisation du format json.")
        output_format = "json"
    
    # Créer le prompt d'évaluateur pour cette matière
    evaluateur_prompt = creer_prompt_evaluateur(matiere, output_format)
    
    # Créer un système RAG avec le prompt d'évaluateur
    retrieval_chain = setup_rag_system(
        index_name=index_name, 
        embeddings=embeddings,
        matiere=matiere, 
        custom_prompt=evaluateur_prompt,
        output_format=output_format
    )
    
    print(f"\n\n--- Évaluation pour {matiere}, question: '{question}' ---")
    
    # Exécuter l'évaluation
    response = retrieval_chain.invoke({
        "input": "Évaluer la réponse de l'étudiant", 
        "question": question,
        "student_response": student_response
    })
    
    # Dossier de sortie pour la session
    output_dir = None
    if save_output:
        output_dir = creer_dossier_sortie()
    
    # Traiter la réponse JSON
    try:
        import json
        import re
        
        # Extraire le JSON de la réponse
        json_answer = response['answer']
        
        # Rechercher le JSON entre les délimiteurs ```json et ```
        json_pattern = r"```(?:json)?(.*?)```"
        match = re.search(json_pattern, json_answer, re.DOTALL)
        
        if match:
            # Extraire le contenu JSON
            json_str = match.group(1).strip()
            evaluation_json = json.loads(json_str)
        else:
            # Essayer de charger directement la réponse comme JSON
            json_str = json_answer.strip()
            evaluation_json = json.loads(json_str)
        
        # Vérifier si des documents d'examen sont présents dans le contexte
        has_exam_docs = any(
            doc.metadata.get("is_exam", False) for doc in response.get("context", [])
        )
        
        # Si le champ "basé_sur_examen" n'est pas déjà défini, l'ajouter
        if "basé_sur_examen" not in evaluation_json:
            evaluation_json["basé_sur_examen"] = has_exam_docs
        
        # Ajouter les sources à la réponse JSON
        sources = []
        for i, doc in enumerate(response["context"]):
            source_entry = {
                "document": i + 1,
                "source": doc.metadata.get('source', 'Source inconnue'),
                "is_exam": doc.metadata.get('is_exam', False)
            }
            
            # Ajouter la section si elle existe
            if "Header 2" in doc.metadata:
                source_entry["section"] = doc.metadata["Header 2"]
            elif "Header 3" in doc.metadata:
                source_entry["section"] = doc.metadata["Header 3"]
            
            sources.append(source_entry)
        
        # Ajouter les sources aux métadonnées
        if "metadonnees" not in evaluation_json:
            evaluation_json["metadonnees"] = {}
        evaluation_json["metadonnees"]["sources"] = sources
        
        # Ajouter des métadonnées supplémentaires
        evaluation_json["metadonnees"]["question"] = question
        evaluation_json["metadonnees"]["date_evaluation"] = datetime.now().isoformat()
        
        # Mettre à jour la réponse
        formatted_json = json.dumps(evaluation_json, ensure_ascii=False, indent=2)
        
        # Stocker le JSON complet dans la réponse
        response['answer'] = formatted_json
        
        # Sauvegarder le résultat si demandé
        if save_output:
            prefix = "evaluation"
            sauvegarder_json(formatted_json, prefix, matiere, output_dir)
        
        print("\nRésultat de l'évaluation (format JSON):")
        print(formatted_json)
    
    except Exception as e:
        print(f"Erreur lors du traitement du JSON d'évaluation: {e}")
    
    return response

# -----------------------------------------
# Fonction principale et utilitaires
# -----------------------------------------
def afficher_matieres_disponibles():
    """
    Affiche la liste des matières disponibles dans le dossier des cours.
    
    Returns:
        list: Liste des matières disponibles
    """
    if not os.path.exists(COURS_DIR):
        print("Le dossier des cours n'existe pas encore.")
        return []
    
    matieres = []
    for item in os.listdir(COURS_DIR):
        item_path = os.path.join(COURS_DIR, item)
        if os.path.isdir(item_path) and not item.startswith('.'):
            matieres.append(item)
    
    if matieres:
        print("Matières disponibles:")
        for matiere in matieres:
            print(f"- {matiere}")
    else:
        print("Aucune matière n'a été trouvée dans le dossier des cours.")
    
    return matieres

def main():
    """
    Fonction principale qui permet d'utiliser le système interactivement.
    """
    print("=== Système de gestion de cours et génération de questions ===")
    
    # Vérifier que l'extraction PDF fonctionne
    print("\n=== Vérification de PyPDF2 ===")
    if PYPDF2_AVAILABLE:
        try:
            # Créer un petit PDF en mémoire pour tester
            print("Test de PyPDF2 avec un document en mémoire...")
            from io import BytesIO
            
            # Créer un PDF minimaliste en mémoire pour tester
            test_successful = False
            try:
                # PyPDF2 3.0+ utilise PdfWriter pour la création
                from PyPDF2 import PdfWriter
                pdf_bytes = BytesIO()
                writer = PdfWriter()
                page = writer.add_blank_page(width=500, height=500)
                writer.add_page(page)
                writer.write(pdf_bytes)
                pdf_bytes.seek(0)
                
                # Lire le PDF créé
                reader = PyPDF2.PdfReader(pdf_bytes)
                print(f"PDF de test créé avec {len(reader.pages)} page(s)")
                test_successful = True
            except Exception as e:
                print(f"Erreur lors du test de création de PDF: {e}")
            
            # Si le test avec PDF en mémoire échoue, essayer avec un fichier réel
            if not test_successful:
                # Chercher un fichier PDF pour tester
                pdf_files = []
                for root, dirs, files in os.walk(COURS_DIR):
                    for file in files:
                        if file.lower().endswith('.pdf'):
                            pdf_files.append(os.path.join(root, file))
                
                if pdf_files:
                    test_file = pdf_files[0]
                    print(f"Test avec le fichier existant: {test_file}")
                    with open(test_file, 'rb') as f:
                        reader = PyPDF2.PdfReader(f)
                        num_pages = len(reader.pages)
                        print(f"Extraction réussie! Le fichier contient {num_pages} pages.")
                        test_successful = True
                else:
                    print("Aucun fichier PDF trouvé pour tester l'extraction.")
            
            if test_successful:
                print("✅ Vérification PyPDF2 terminée avec succès!")
            else:
                print("⚠️ Impossible de valider complètement la fonctionnalité PyPDF2.")
        except Exception as e:
            print(f"❌ ERREUR: Problème avec PyPDF2: {e}")
            print("Les extractions PDF peuvent échouer pendant l'exécution!")
    else:
        print("❌ PyPDF2 n'est pas disponible. Les extractions PDF seront limitées.")
        print("Installez PyPDF2 avec: pip install PyPDF2==3.0.1")
    
    # Initialiser la structure des dossiers
    initialiser_structure_dossiers()
    
    # Initialiser Pinecone
    pc, index_name, spec = initialize_pinecone()
    
    # Configurer le modèle d'embedding
    embeddings = setup_embeddings()
    
    # Créer ou récupérer l'index
    index = create_or_get_index(pc, index_name, embeddings, spec)
    
    # Afficher les matières disponibles
    matieres = afficher_matieres_disponibles()
    
    # Menu interactif
    while True:
        print("\nQue souhaitez-vous faire?")
        print("1. Mettre à jour les documents d'une matière")
        print("2. Générer une question de réflexion (texte)")
        print("3. Générer une question de réflexion (JSON)")
        print("4. Poser une question sur une matière")
        print("5. Poser une question (réponse JSON)")
        print("6. Évaluer la réponse d'un étudiant")
        print("7. Quitter")
        
        choix = input("Votre choix (1-7): ")
        
        if choix == "1":
            # Mettre à jour une matière
            if not matieres:
                print("Aucune matière disponible. Veuillez ajouter des documents dans le dossier 'cours'.")
                continue
                
            matiere = input(f"Entrez le nom de la matière à mettre à jour ({', '.join(matieres)}): ").upper()
            if matiere not in matieres:
                print(f"Matière '{matiere}' non trouvée.")
                continue
                
            mettre_a_jour_matiere(pc, index_name, embeddings, matiere)
            
        elif choix == "2":
            # Générer une question de réflexion (texte)
            if not matieres:
                print("Aucune matière disponible.")
                continue
                
            matiere = input(f"Entrez le nom de la matière ({', '.join(matieres)}): ").upper()
            if matiere not in matieres:
                print(f"Matière '{matiere}' non trouvée.")
                continue
                
            concept = input("Entrez le concept sur lequel générer une question: ")
            
            try:
                question = generer_question_reflexion(index_name, embeddings, matiere, concept, "text", save_output=False)
                print(f"\nQuestion générée: {question}")
            except Exception as e:
                print(f"Erreur lors de la génération de la question: {e}")
                
        elif choix == "3":
            # Générer une question de réflexion (JSON)
            if not matieres:
                print("Aucune matière disponible.")
                continue
                
            matiere = input(f"Entrez le nom de la matière ({', '.join(matieres)}): ").upper()
            if matiere not in matieres:
                print(f"Matière '{matiere}' non trouvée.")
                continue
                
            concept = input("Entrez le concept sur lequel générer une question: ")
            
            try:
                question_json = generer_question_reflexion(index_name, embeddings, matiere, concept, "json", save_output=True)
                print(f"\nQuestion JSON générée et sauvegardée.\n{question_json}")
            except Exception as e:
                print(f"Erreur lors de la génération de la question JSON: {e}")
            
        elif choix == "4":
            # Poser une question
            if not matieres:
                print("Aucune matière disponible.")
                continue
                
            matiere = input(f"Entrez le nom de la matière ({', '.join(matieres)}): ").upper()
            if matiere not in matieres:
                print(f"Matière '{matiere}' non trouvée.")
                continue
                
            question = input("Entrez votre question: ")
            
            try:
                interroger_matiere(index_name, embeddings, matiere, question, save_output=False)
            except Exception as e:
                print(f"Erreur lors de la recherche: {e}")
                
        elif choix == "5":
            # Poser une question avec réponse JSON
            if not matieres:
                print("Aucune matière disponible.")
                continue
                
            matiere = input(f"Entrez le nom de la matière ({', '.join(matieres)}): ").upper()
            if matiere not in matieres:
                print(f"Matière '{matiere}' non trouvée.")
                continue
                
            question = input("Entrez votre question: ")
            
            try:
                interroger_matiere(index_name, embeddings, matiere, question, output_format="json", save_output=True)
            except Exception as e:
                print(f"Erreur lors de la recherche: {e}")
                
        elif choix == "6":
            # Évaluer la réponse d'un étudiant
            if not matieres:
                print("Aucune matière disponible.")
                continue
                
            matiere = input(f"Entrez le nom de la matière ({', '.join(matieres)}): ").upper()
            if matiere not in matieres:
                print(f"Matière '{matiere}' non trouvée.")
                continue
                
            question = input("Entrez la question posée à l'étudiant: ")
            print("\nEntrez la réponse de l'étudiant (terminez par une ligne contenant uniquement 'FIN'):")
            
            lines = []
            while True:
                line = input()
                if line.strip() == "FIN":
                    break
                lines.append(line)
            
            student_response = "\n".join(lines)
            
            if not student_response.strip():
                print("La réponse de l'étudiant ne peut pas être vide.")
                continue
                
            try:
                evaluation = evaluer_reponse_etudiant(
                    index_name, 
                    embeddings, 
                    matiere, 
                    question, 
                    student_response, 
                    output_format="json", 
                    save_output=True
                )
                print("\nÉvaluation terminée et sauvegardée.")
            except Exception as e:
                print(f"Erreur lors de l'évaluation: {e}")
            
        elif choix == "7":
            # Quitter
            print("Au revoir!")
            break
            
        else:
            print("Choix non valide. Veuillez entrer un nombre entre 1 et 7.")

if __name__ == "__main__":
    main() 